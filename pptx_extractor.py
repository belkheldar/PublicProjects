#!/usr/bin/env python3
"""
pptx_extractor.py
-----------------
Extracts a PowerPoint (.pptx) file into a structured JSON format
designed to be fed into an AI for regenerating presentations.

Usage:
    python pptx_extractor.py <input.pptx> [output.json]

Output JSON structure:
    {
        "meta": { title, author, slide_count, slide_width_px, slide_height_px },
        "theme": { colors, fonts },
        "slides": [
            {
                "index": 1,
                "layout_name": "Title Slide",
                "elements": [
                    {
                        "type": "text|image|chart|table|shape",
                        "role": "title|subtitle|body|label|...",
                        "position": { x_pct, y_pct, w_pct, h_pct },
                        "text": "...",          # for text elements
                        "paragraphs": [...],    # rich text breakdown
                        "style": { font, size, bold, italic, color, align },
                        "image_path": "...",    # for image elements
                        "image_description": "", # placeholder for AI vision pass
                        "chart_type": "...",    # for chart elements
                        "chart_data": {...},
                        "table_data": [...],    # for table elements
                    }
                ],
                "notes": "speaker notes text",
                "background": { type, color }
            }
        ]
    }
"""

import json
import sys
import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64


# ── helpers ──────────────────────────────────────────────────────────────────

def emu_to_pct(value_emu: int, dimension_emu: int) -> float:
    """Convert EMU position to percentage of slide dimension."""
    if dimension_emu == 0:
        return 0.0
    return round((value_emu / dimension_emu) * 100, 2)


def rgb_to_hex(rgb) -> str | None:
    if rgb is None:
        return None
    return "#{:02X}{:02X}{:02X}".format(rgb.r, rgb.g, rgb.b)


def safe_str(value) -> str | None:
    try:
        return str(value) if value is not None else None
    except Exception:
        return None


def align_name(align) -> str:
    mapping = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
        PP_ALIGN.DISTRIBUTE: "distribute",
    }
    return mapping.get(align, "left")


def pt_size(font_size_emu) -> float | None:
    """Convert font size from EMU (hundredths of a point) to points."""
    if font_size_emu is None:
        return None
    return round(font_size_emu / 12700, 1)


# ── element extractors ───────────────────────────────────────────────────────

def extract_run_style(run) -> dict:
    font = run.font
    return {
        "font_name": font.name,
        "font_size_pt": pt_size(font.size),
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "color_hex": rgb_to_hex(font.color.rgb) if font.color and font.color.type else None,
    }


def extract_paragraph(para) -> dict:
    runs = []
    for run in para.runs:
        runs.append({
            "text": run.text,
            "style": extract_run_style(run),
        })

    # paragraph-level style
    pf = para.font if hasattr(para, "font") else None
    return {
        "text": para.text,
        "alignment": align_name(para.alignment) if para.alignment else None,
        "level": para.level,
        "runs": runs,
    }


def infer_role(shape) -> str:
    """Guess semantic role from placeholder type or shape name."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    if shape.is_placeholder:
        ph = shape.placeholder_format
        role_map = {
            PP_PLACEHOLDER.TITLE: "title",
            PP_PLACEHOLDER.CENTER_TITLE: "title",
            PP_PLACEHOLDER.SUBTITLE: "subtitle",
            PP_PLACEHOLDER.BODY: "body",
            PP_PLACEHOLDER.OBJECT: "object",
            PP_PLACEHOLDER.PICTURE: "image",
            PP_PLACEHOLDER.CHART: "chart",
            PP_PLACEHOLDER.TABLE: "table",
            PP_PLACEHOLDER.DATE: "date",
            PP_PLACEHOLDER.FOOTER: "footer",
            PP_PLACEHOLDER.SLIDE_NUMBER: "slide_number",
        }
        return role_map.get(ph.type, "placeholder")
    name = (shape.name or "").lower()
    for keyword in ["title", "subtitle", "body", "image", "logo", "footer", "header", "chart", "table"]:
        if keyword in name:
            return keyword
    return "content"


def extract_text_element(shape, slide_w, slide_h) -> dict:
    tf = shape.text_frame
    paragraphs = [extract_paragraph(p) for p in tf.paragraphs]
    full_text = "\n".join(p["text"] for p in paragraphs if p["text"])

    # dominant style: from first non-empty run
    dominant_style = {}
    for para in paragraphs:
        for run in para.get("runs", []):
            if run.get("text"):
                dominant_style = run.get("style", {})
                break
        if dominant_style:
            break

    return {
        "type": "text",
        "role": infer_role(shape),
        "name": shape.name,
        "position": {
            "x_pct": emu_to_pct(shape.left, slide_w),
            "y_pct": emu_to_pct(shape.top, slide_h),
            "w_pct": emu_to_pct(shape.width, slide_w),
            "h_pct": emu_to_pct(shape.height, slide_h),
        },
        "text": full_text,
        "paragraphs": paragraphs,
        "dominant_style": dominant_style,
        "word_wrap": tf.word_wrap,
    }


def extract_image_element(shape, slide_w, slide_h, image_dir: Path, slide_idx: int) -> dict:
    image = shape.image
    ext = image.ext or "png"
    filename = f"slide{slide_idx:02d}_{shape.shape_id}.{ext}"
    image_path = image_dir / filename
    image_path.write_bytes(image.blob)

    return {
        "type": "image",
        "role": infer_role(shape),
        "name": shape.name,
        "position": {
            "x_pct": emu_to_pct(shape.left, slide_w),
            "y_pct": emu_to_pct(shape.top, slide_h),
            "w_pct": emu_to_pct(shape.width, slide_w),
            "h_pct": emu_to_pct(shape.height, slide_h),
        },
        "image_path": str(image_path),
        "content_type": image.content_type,
        "image_description": "",  # fill via AI vision pass
    }


def extract_chart_element(shape, slide_w, slide_h) -> dict:
    chart = shape.chart
    chart_type = safe_str(chart.chart_type)

    series_data = []
    try:
        for series in chart.series:
            values = []
            try:
                values = [v for v in series.values]
            except Exception:
                pass
            series_data.append({
                "name": safe_str(series.name),
                "values": values,
            })
    except Exception:
        pass

    categories = []
    try:
        cats = chart.plots[0].series[0].data_labels if chart.plots else None
        # Try to get category names via the plot
        plot = chart.plots[0]
        if hasattr(plot, "series") and plot.series:
            s = plot.series[0]
            if hasattr(s, "data_labels"):
                pass
        # Simpler approach: iterate chart data
        for plot in chart.plots:
            for series in plot.series:
                if hasattr(series, "values"):
                    break
    except Exception:
        pass

    return {
        "type": "chart",
        "role": infer_role(shape),
        "name": shape.name,
        "position": {
            "x_pct": emu_to_pct(shape.left, slide_w),
            "y_pct": emu_to_pct(shape.top, slide_h),
            "w_pct": emu_to_pct(shape.width, slide_w),
            "h_pct": emu_to_pct(shape.height, slide_h),
        },
        "chart_type": chart_type,
        "chart_title": safe_str(chart.chart_title.text_frame.text) if chart.has_title else None,
        "series": series_data,
        "image_description": "",  # fill via AI vision pass
    }


def extract_table_element(shape, slide_w, slide_h) -> dict:
    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append({
                "text": cell.text,
                "paragraphs": [extract_paragraph(p) for p in cell.text_frame.paragraphs],
            })
        rows.append(cells)

    return {
        "type": "table",
        "role": "table",
        "name": shape.name,
        "position": {
            "x_pct": emu_to_pct(shape.left, slide_w),
            "y_pct": emu_to_pct(shape.top, slide_h),
            "w_pct": emu_to_pct(shape.width, slide_w),
            "h_pct": emu_to_pct(shape.height, slide_h),
        },
        "row_count": len(table.rows),
        "col_count": len(table.columns),
        "rows": rows,
    }


def extract_shape_element(shape, slide_w, slide_h) -> dict:
    el = {
        "type": "shape",
        "role": infer_role(shape),
        "name": shape.name,
        "position": {
            "x_pct": emu_to_pct(shape.left, slide_w),
            "y_pct": emu_to_pct(shape.top, slide_h),
            "w_pct": emu_to_pct(shape.width, slide_w),
            "h_pct": emu_to_pct(shape.height, slide_h),
        },
        "shape_type": safe_str(shape.shape_type),
    }
    if shape.has_text_frame:
        tf = shape.text_frame
        el["text"] = "\n".join(p.text for p in tf.paragraphs)
        el["paragraphs"] = [extract_paragraph(p) for p in tf.paragraphs]

    # fill color
    try:
        fill = shape.fill
        if fill.type is not None:
            el["fill_color"] = rgb_to_hex(fill.fore_color.rgb) if fill.fore_color else None
    except Exception:
        pass

    return el


# ── background ───────────────────────────────────────────────────────────────

def extract_background(slide) -> dict:
    bg = {"type": "none", "color": None}
    try:
        fill = slide.background.fill
        if fill.type is not None:
            bg["type"] = safe_str(fill.type)
            try:
                bg["color"] = rgb_to_hex(fill.fore_color.rgb)
            except Exception:
                pass
    except Exception:
        pass
    return bg


# ── theme ─────────────────────────────────────────────────────────────────────

def extract_theme(prs) -> dict:
    theme = {"colors": [], "fonts": {"major": None, "minor": None}}
    try:
        slide_master = prs.slide_master
        # Theme colors from XML
        from lxml import etree
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        theme_elem = slide_master._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme")
        if theme_elem is not None:
            for clr in theme_elem.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"):
                val = clr.get("val")
                if val:
                    theme["colors"].append(f"#{val.upper()}")
        theme["colors"] = list(dict.fromkeys(theme["colors"]))[:12]  # dedupe, limit

        # Fonts
        font_scheme = theme_elem.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}fontScheme") if theme_elem is not None else None
        if font_scheme is not None:
            major = font_scheme.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}majorFont/{http://schemas.openxmlformats.org/drawingml/2006/main}latin")
            minor = font_scheme.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}minorFont/{http://schemas.openxmlformats.org/drawingml/2006/main}latin")
            theme["fonts"]["major"] = major.get("typeface") if major is not None else None
            theme["fonts"]["minor"] = minor.get("typeface") if minor is not None else None
    except Exception:
        pass
    return theme


# ── main extractor ────────────────────────────────────────────────────────────

def extract_pptx(pptx_path: str, output_path: str | None = None) -> dict:
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")

    prs = Presentation(str(pptx_path))

    # Output paths
    if output_path is None:
        output_path = pptx_path.with_suffix(".json")
    else:
        output_path = Path(output_path)

    image_dir = output_path.parent / f"{output_path.stem}_images"
    image_dir.mkdir(parents=True, exist_ok=True)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Meta
    props = prs.core_properties
    meta = {
        "source_file": pptx_path.name,
        "title": props.title or "",
        "author": props.author or "",
        "slide_count": len(prs.slides),
        "slide_width_emu": slide_w,
        "slide_height_emu": slide_h,
        "slide_width_px": round(slide_w / 9144),   # 1px ≈ 9144 EMU at 96dpi
        "slide_height_px": round(slide_h / 9144),
        "aspect_ratio": f"{round(slide_w/slide_h, 4)}",
    }

    theme = extract_theme(prs)

    slides = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"

        elements = []
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Flatten group: iterate sub-shapes
                    for sub in shape.shapes:
                        try:
                            el = _extract_single_shape(sub, slide_w, slide_h, image_dir, slide_idx)
                            if el:
                                elements.append(el)
                        except Exception as e:
                            elements.append({"type": "error", "name": getattr(sub, "name", "?"), "error": str(e)})
                else:
                    el = _extract_single_shape(shape, slide_w, slide_h, image_dir, slide_idx)
                    if el:
                        elements.append(el)
            except Exception as e:
                elements.append({"type": "error", "name": getattr(shape, "name", "?"), "error": str(e)})

        # Speaker notes
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        slides.append({
            "index": slide_idx,
            "layout_name": layout_name,
            "background": extract_background(slide),
            "elements": elements,
            "notes": notes_text,
        })

    result = {
        "meta": meta,
        "theme": theme,
        "slides": slides,
    }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2, ensure_ascii=False, default=str)

    print(f"✓ Extracted {len(slides)} slides → {output_path}")
    if any(image_dir.iterdir()):
        print(f"✓ Images saved → {image_dir}/")
    return result


def _extract_single_shape(shape, slide_w, slide_h, image_dir, slide_idx) -> dict | None:
    """Route a single shape to the right extractor."""
    # Skip invisible/zero-size shapes
    if shape.width == 0 or shape.height == 0:
        return None

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return extract_image_element(shape, slide_w, slide_h, image_dir, slide_idx)

    if shape.has_chart:
        return extract_chart_element(shape, slide_w, slide_h)

    if shape.has_table:
        return extract_table_element(shape, slide_w, slide_h)

    if shape.has_text_frame:
        return extract_text_element(shape, slide_w, slide_h)

    return extract_shape_element(shape, slide_w, slide_h)


# ── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pptx_extractor.py <input.pptx> [output.json]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    result = extract_pptx(input_file, output_file)

    # Print a quick summary
    print("\n── Summary ──")
    for slide in result["slides"]:
        title_el = next((e for e in slide["elements"] if e.get("role") == "title"), None)
        title = title_el["text"] if title_el else "(no title)"
        elem_types = [e["type"] for e in slide["elements"]]
        type_counts = {}
        for t in elem_types:
            type_counts[t] = type_counts.get(t, 0) + 1
        counts = ", ".join(f"{v}x {k}" for k, v in type_counts.items())
        print(f"  Slide {slide['index']:2d} [{slide['layout_name']}]: {title[:50]!r} — {counts}")
